import os
import time
import zipfile
import qrcode 
import math
import json
import logging
from io import BytesIO 
from flask import Flask, render_template, request, send_file, jsonify
from werkzeug.utils import secure_filename
from threading import Thread
from PyPDF2 import PdfReader, PdfWriter, PdfMerger
from pdf2image import convert_from_path
from PIL import Image 
import io
import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Inches
from flask_cors import CORS
from config import config
from dotenv import load_dotenv

try:
    import google.generativeai as genai
    GENAI_IMPORT_ERROR = None
except Exception as _genai_err:
    genai = None
    GENAI_IMPORT_ERROR = str(_genai_err)

try:
    import pythoncom
    import win32com.client
    HAS_WIN32COM = True
except ImportError:
    pythoncom = None
    win32com = None
    HAS_WIN32COM = False

# Load environment variables
load_dotenv()

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# --- SECURE AI CONFIGURATION ---
GENAI_AVAILABLE = True
if GENAI_IMPORT_ERROR:
    GENAI_AVAILABLE = False
    logger.warning("google.generativeai import unavailable: %s", GENAI_IMPORT_ERROR)

if not config.GOOGLE_API_KEY:
    GENAI_AVAILABLE = False
    logger.warning("GOOGLE_API_KEY not found; AI features will be disabled")

if GENAI_AVAILABLE:
    try:
        genai.configure(api_key=config.GOOGLE_API_KEY)
    except Exception as _cfg_err:
        GENAI_AVAILABLE = False
        logger.warning("Failed to configure Gemini client: %s", _cfg_err)

AI_MODEL_ID = config.AI_MODEL_ID

# Initialize Flask app
app = Flask(__name__)
app.config.from_object(config)

# Configure CORS with allowed origins - Apply to all routes
CORS(app, 
     origins=config.CORS_ORIGINS,
     allow_headers=['Content-Type', 'Authorization'],
     methods=['GET', 'POST', 'PUT', 'DELETE', 'OPTIONS'],
     supports_credentials=True,
     max_age=3600)

# Create required folders
UPLOAD_FOLDER = config.UPLOAD_FOLDER
VAULT_FOLDER = config.VAULT_FOLDER
DATA_FILE = os.path.join(UPLOAD_FOLDER, 'student_data.json')
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs(VAULT_FOLDER, exist_ok=True)

# --- SECURITY HEADERS MIDDLEWARE ---
@app.after_request
def add_security_headers(response):
    """Add security headers to all responses"""
    response.headers['X-Content-Type-Options'] = 'nosniff'
    response.headers['X-Frame-Options'] = 'SAMEORIGIN'
    response.headers['X-XSS-Protection'] = '1; mode=block'
    response.headers['Strict-Transport-Security'] = 'max-age=31536000; includeSubDomains'
    return response

# --- NEW DATA STORAGE HELPERS ---
def load_student_data():
    if os.path.exists(DATA_FILE):
        try:
            with open(DATA_FILE, 'r') as f:
                return json.load(f)
        except:
            pass
    # Updated to include 'vault' list
    return {"schedule": [], "todo": [], "assignments": [], "study": [], "vault": []}

def save_student_data(data):
    with open(DATA_FILE, 'w') as f:
        json.dump(data, f)

# --- TEST ENDPOINT FOR DEBUGGING ---
@app.route('/test-connection', methods=['GET', 'POST'])
def test_connection():
    """Test endpoint to verify backend is responding"""
    return jsonify({
        "status": "ok",
        "message": "Backend is responding correctly",
        "timestamp": time.time()
    })

# --- LOGIC FUNCTIONS (STAYED THE SAME) ---

def merge_pdfs_logic(input_paths):
    merger = PdfMerger()
    for path in input_paths:
        merger.append(path)
    output_filename = f"merged_{int(time.time())}.pdf"
    output_path = os.path.join(UPLOAD_FOLDER, output_filename)
    with open(output_path, "wb") as f:
        merger.write(f)
    merger.close()
    return output_path

def split_pdf_logic(input_path, start, end):
    reader = PdfReader(input_path)
    total_pages = len(reader.pages)
    if not start or start < 1: start = 1
    if not end or end > total_pages: end = total_pages
    
    writer = PdfWriter()
    for i in range(start - 1, end):
        writer.add_page(reader.pages[i])
    
    output_filename = f"extracted_{start}_to_{end}.pdf"
    output_path = os.path.join(UPLOAD_FOLDER, output_filename)
    with open(output_path, "wb") as f:
        writer.write(f)
    return output_path

def protect_pdf_logic(input_path, password):
    reader = PdfReader(input_path)
    writer = PdfWriter()
    for page in reader.pages:
        writer.add_page(page)
    
    writer.encrypt(password)
    
    output_filename = f"protected_{int(time.time())}.pdf"
    output_path = os.path.join(UPLOAD_FOLDER, output_filename)
    with open(output_path, "wb") as f:
        writer.write(f)
    return output_path

def pdf_to_img_logic(input_path, page_str):
    pages_to_convert = []
    if page_str and page_str.strip():
        try:
            pages_to_convert = [int(p.strip()) for p in page_str.split(',')]
        except Exception as e:
            logger.warning(f"Error parsing page numbers: {e}")

    images = convert_from_path(
        input_path, 
        poppler_path=config.POPPLER_PATH
    )
    
    saved_files = []
    for i, image in enumerate(images):
        page_num = i + 1
        if not pages_to_convert or page_num in pages_to_convert:
            img_name = f"page_{page_num}_{int(time.time())}.jpg"
            img_path = os.path.join(UPLOAD_FOLDER, img_name)
            image.save(img_path, 'JPEG')
            saved_files.append(img_path)
            
    return saved_files

def convert_to_pdf(in_path, out_path):
    if not HAS_WIN32COM:
        raise RuntimeError("Word to PDF conversion requires Windows with pywin32 installed")

    pythoncom.CoInitialize()
    try:
        word = win32com.client.Dispatch("Word.Application")
        word.Visible = False
        doc = word.Documents.Open(in_path)
        doc.SaveAs(out_path, FileFormat=17)
        doc.Close()
        word.Quit()
    finally:
        pythoncom.CoUninitialize()

# --- ROUTES ---

@app.route('/')
def index():
    return render_template('index.html')

@app.route('/convert', methods=['POST'])
def handle_conversion():
    files = request.files.getlist('file')
    if not files or files[0].filename == '':
        return "No files", 400

    if not HAS_WIN32COM:
        return (
            "Word to PDF conversion is not available on this hosting environment. "
            "Use local Windows deployment for this feature.",
            501,
        )

    results = []
    failed_files = []
    for file in files:
        fname = secure_filename(file.filename)
        in_p = os.path.join(UPLOAD_FOLDER, fname)
        file.save(in_p)
        out_fname = fname.rsplit('.', 1)[0] + '.pdf'
        out_p = os.path.join(UPLOAD_FOLDER, out_fname)
        try:
            convert_to_pdf(in_p, out_p)
            results.append(out_p)
        except Exception as e:
            logger.warning("Word conversion failed for %s: %s", fname, e)
            failed_files.append(fname)

    if not results:
        return (
            "Could not convert the uploaded Word file(s) on this server. "
            "This feature requires Microsoft Word automation on Windows.",
            501,
        )

    if len(results) > 1:
        zip_p = os.path.join(UPLOAD_FOLDER, f"converted_{int(time.time())}.zip")
        with zipfile.ZipFile(zip_p, 'w') as z:
            for f in results: z.write(f, os.path.basename(f))
        return send_file(zip_p, as_attachment=True)
    return send_file(results[0], as_attachment=True)

@app.route('/split', methods=['POST'])
def handle_split():
    file = request.files.get('file')
    start = request.form.get('start_page', type=int)
    end = request.form.get('end_page', type=int)
    if not file: return "No file", 400
    fname = secure_filename(file.filename)
    in_path = os.path.join(UPLOAD_FOLDER, fname)
    file.save(in_path)
    try:
        extracted_pdf_path = split_pdf_logic(in_path, start, end)
        return send_file(extracted_pdf_path, as_attachment=True)
    except Exception as e:
        return f"Split failed: {str(e)}", 500

@app.route('/merge', methods=['POST'])
def handle_merge():
    files = request.files.getlist('file')
    if not files or files[0].filename == '': return "No files selected", 400
    file_paths = []
    for file in files:
        fname = secure_filename(file.filename)
        file_p = os.path.join(UPLOAD_FOLDER, fname)
        file.save(file_p)
        file_paths.append(file_p)
    try:
        merged_pdf_path = merge_pdfs_logic(file_paths)
        return send_file(merged_pdf_path, as_attachment=True)
    except Exception as e:
        return f"Merge failed: {str(e)}", 500

@app.route('/protect', methods=['POST'])
def handle_protect():
    file = request.files.get('file')
    password = request.form.get('pdf_password')
    if not file or not password: return "File and password required", 400
    fname = secure_filename(file.filename)
    in_p = os.path.join(UPLOAD_FOLDER, fname)
    file.save(in_p)
    try:
        protected_p = protect_pdf_logic(in_p, password)
        return send_file(protected_p, as_attachment=True)
    except Exception as e:
        return f"Protection failed: {str(e)}", 500

@app.route('/pdf-to-img', methods=['POST'])
def handle_pdf_to_img():
    file = request.files.get('file')
    page_str = request.form.get('pages', "")
    if not file: return "No file", 400
    fname = secure_filename(file.filename)
    in_p = os.path.join(UPLOAD_FOLDER, fname)
    file.save(in_p)
    try:
        img_paths = pdf_to_img_logic(in_p, page_str)
        if len(img_paths) > 1:
            zip_filename = f"images_{int(time.time())}.zip"
            zip_p = os.path.join(UPLOAD_FOLDER, zip_filename)
            with zipfile.ZipFile(zip_p, 'w') as z:
                for img in img_paths:
                    z.write(img, os.path.basename(img))
            return send_file(zip_p, as_attachment=True)
        return send_file(img_paths[0], as_attachment=True)
    except Exception as e:
        return f"Image conversion failed: {str(e)}", 500

@app.route('/generate-qr', methods=['POST'])
def handle_qr():
    data = request.form.get('qr_data')
    if not data: return "No data provided", 400
    
    qr = qrcode.QRCode(version=1, box_size=10, border=5)
    qr.add_data(data)
    qr.make(fit=True)
    img = qr.make_image(fill_color="black", back_color="white")
    
    img_io = BytesIO()
    img.save(img_io, 'PNG')
    img_io.seek(0)
    return send_file(img_io, mimetype='image/png', as_attachment=True, download_name='qrcode.png')

@app.route('/img-to-pdf', methods=['POST'])
def handle_img_to_pdf():
    files = request.files.getlist('file')
    if not files or files[0].filename == '': return "No files", 400
    
    image_list = []
    try:
        for file in files:
            fname = secure_filename(file.filename)
            in_p = os.path.join(UPLOAD_FOLDER, fname)
            file.save(in_p)
            
            img = Image.open(in_p)
            if img.mode != 'RGB':
                img = img.convert('RGB')
            image_list.append(img)
            
        if not image_list: return "No valid images found", 400
        
        output_filename = f"converted_{int(time.time())}.pdf"
        output_path = os.path.join(UPLOAD_FOLDER, output_filename)
        image_list[0].save(output_path, save_all=True, append_images=image_list[1:])
        
        return send_file(output_path, as_attachment=True)
    except Exception as e:
        return f"Image to PDF failed: {str(e)}", 500
    
    # ... (existing img-to-pdf code) ...

# --- ADD THIS TO YOUR app.py ---

@app.route('/convert-pdf-to-ppt', methods=['POST'])
def convert_pdf_to_ppt():
    try:
        if 'file' not in request.files:
            return jsonify({"error": "No file uploaded"}), 400
        
        pdf_file = request.files['file']
        
        # 1. Open PDF directly from memory
        pdf_stream = pdf_file.read()
        pdf_document = fitz.open(stream=pdf_stream, filetype="pdf")
        
        # 2. Create PowerPoint object
        prs = Presentation()
        
        for page_num in range(len(pdf_document)):
            page = pdf_document.load_page(page_num)
            
            # Render page to high-res image (300 DPI approx)
            pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
            img_data = io.BytesIO(pix.tobytes("png"))
            
            # Add a blank slide (layout 6 is usually blank)
            slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(slide_layout)
            
            # Stretch image to fill the whole slide
            slide.shapes.add_picture(img_data, 0, 0, width=prs.slide_width, height=prs.slide_height)
        
        # 3. Save PPT to memory and send back to browser
        ppt_output = io.BytesIO()
        prs.save(ppt_output)
        ppt_output.seek(0)
        
        return send_file(
            ppt_output,
            as_attachment=True,
            download_name="AcademiaPro_Converted.pptx",
            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation'
        )

    except Exception as e:
        print(f"Conversion Error: {str(e)}")
        return jsonify({"error": str(e)}), 500

@app.route('/calculate-attendance', methods=['POST'])
def handle_attendance():
    try:
        held = int(request.form.get('classes_held', 0))
        attended = int(request.form.get('classes_attended', 0))
        required_raw = request.form.get('min_percentage', '75')
        required = float(required_raw) if required_raw else 75.0
        
        current_pct = (attended / held) * 100 if held > 0 else 0
        
        needed = 0
        if current_pct < required:
            needed = math.ceil((required * held - 100 * attended) / (100 - required))
        
        status = "Safe" if current_pct >= required else "Short Attendance"
        color = "#059669" if status == "Safe" else "#df481a"
        
        return f"""
        <div style="font-family: sans-serif; text-align: center; padding: 50px;">
            <h1 style="color: {color};">{status}</h1>
            <p style="font-size: 1.2rem;">Your current attendance is <strong>{current_pct:.2f}%</strong></p>
            {f'<p style="color: #df481a; font-weight: bold;">You need to attend <strong>{needed}</strong> more classes consecutively to hit {required}%.</p>' if needed > 0 else '<p style="color: #059669;">You are already above the required limit!</p>'}
            <br><br>
            <a href="/" style="text-decoration: none; color: white; background: #1e293b; padding: 10px 20px; border-radius: 5px;">Back to Dashboard</a>
        </div>
        """
    except Exception as e:
        return f"Calculation Error: {str(e)}", 400

# --- STUDENT HUB ROUTES ---

@app.route('/get-student-info', methods=['GET'])
def get_info():
    return jsonify(load_student_data())

@app.route('/add-student-info', methods=['POST'])
def add_info():
    data = request.json
    category = data.get('category')
    content = data.get('content')
    
    all_data = load_student_data()
    if category in all_data:
        all_data[category].append(content)
        save_student_data(all_data)
        return jsonify({"status": "success"})
    return jsonify({"status": "error"}), 400

@app.route('/delete-student-info', methods=['POST'])
def delete_info():
    data = request.json
    category = data.get('category')
    index = data.get('index')
    
    all_data = load_student_data()
    if category in all_data and 0 <= index < len(all_data[category]):
        # Permanent file deletion logic if it's from the vault
        if category == 'vault':
            file_to_del = all_data[category][index].split('|')[2]
            if os.path.exists(file_to_del):
                os.remove(file_to_del)
        all_data[category].pop(index)
        save_student_data(all_data)
        return jsonify({"status": "success"})
    return jsonify({"status": "error"}), 400

# --- PERMANENT VAULT STORAGE ROUTE ---
@app.route('/upload-to-vault', methods=['POST'])
def vault_upload():
    file = request.files.get('file')
    custom_name = request.form.get('custom_name', 'Untitled')
    category = request.form.get('category', 'General')
    if not file: return "No file", 400
    
    timestamp = int(time.time())
    fname = f"{timestamp}_{secure_filename(file.filename)}"
    save_path = os.path.join(VAULT_FOLDER, fname)
    file.save(save_path)
    
    entry = f"{custom_name}|{category}|{save_path}|{timestamp}"
    data = load_student_data()
    data['vault'].append(entry)
    save_student_data(data)
    return jsonify({"status": "success"})

@app.route('/download-vault/<int:index>')
def vault_download(index):
    data = load_student_data()
    if 0 <= index < len(data['vault']):
        parts = data['vault'][index].split('|')
        return send_file(parts[2], as_attachment=True, download_name=parts[0] + os.path.splitext(parts[2])[1])
    return "Not found", 404

# --- AI ASSISTANCE ROUTES (STABLE v1 UPDATE) ---

# --- AI ASSISTANCE ROUTES (ENHANCED PROMPT ENGINEERING) ---

@app.route('/ai-academic-help', methods=['POST'])
def ai_academic_help():
    try:
        if not GENAI_AVAILABLE:
            return jsonify({"answer": "AI service is temporarily unavailable on this deployment runtime. Please set Python 3.11.9 in Render and redeploy."}), 503

        data = request.get_json()
        user_query = data.get('query')
        
        # --- ENHANCED SYSTEM INSTRUCTIONS ---
        # This tells the AI how to behave before it reads the student's query.
        system_instruction = (
            "You are the AcademiaPro Academic Tutor. Your goal is to help university students "
            "excel in their studies. Since the user is a university student, "
            "provide detailed, scientifically accurate, yet easy-to-understand explanations. "
            "Use clear headings, bold key terms, and bullet points for lists. "
            "Focus on being supportive, professional, and educationally focused."
        )
        
        # Combine the system instruction with the user's specific query
        full_prompt = f"{system_instruction}\n\nStudent Query: {user_query}"
        
        model = genai.GenerativeModel(AI_MODEL_ID)
        response = model.generate_content(full_prompt)
        return jsonify({"answer": response.text})
        
    except Exception as e:
        error_msg = str(e)
        print(f"--- AI ROUTE ERROR: {error_msg} ---") 
        
        # Handle specific network errors gracefully so the app feels stable
        if "11001" in error_msg or "getaddrinfo" in error_msg:
            return jsonify({"answer": "Network Connection Error: AcademiaPro cannot reach the AI servers. Please check your internet."}), 503
        if "EOF" in error_msg:
            return jsonify({"answer": "Connection lost while generating. Please try asking again."}), 503
            
        return jsonify({"answer": f"AcademiaPro is experiencing a glitch: {error_msg}"}), 500
    # --- INDEPENDENT AI RESEARCH LAB ROUTES ---

@app.route('/ai_lab_upload', methods=['POST'])
def ai_lab_upload():
    """
    Upload documents to AI Research Lab for analysis
    Files are stored permanently and can be analyzed later
    """
    try:
        file = request.files.get('file')
        
        if not file:
            return jsonify({"error": "No file selected"}), 400
        
        filename = secure_filename(file.filename)
        LAB_FOLDER = os.path.join(UPLOAD_FOLDER, 'ai_lab_docs')
        os.makedirs(LAB_FOLDER, exist_ok=True)
        
        # Save file permanently
        file_path = os.path.join(LAB_FOLDER, filename)
        file.save(file_path)
        
        logger.info(f"File uploaded to AI Lab: {filename}")
        
        return jsonify({
            "status": "success",
            "filename": filename,
            "message": f"File '{filename}' uploaded successfully to AI Research Lab. You can now analyze it with AI."
        })
        
    except Exception as e:
        logger.error(f"Upload Error: {str(e)}")
        return jsonify({"error": str(e)}), 500

@app.route('/get-ai-lab-files', methods=['GET'])
def get_ai_lab_files():
    """
    Get list of all files in AI Research Lab
    """
    try:
        LAB_FOLDER = os.path.join(UPLOAD_FOLDER, 'ai_lab_docs')
        
        if not os.path.exists(LAB_FOLDER):
            return jsonify({"files": []})
        
        files = []
        for filename in os.listdir(LAB_FOLDER):
            file_path = os.path.join(LAB_FOLDER, filename)
            if os.path.isfile(file_path):
                file_size = os.path.getsize(file_path)
                file_time = os.path.getmtime(file_path)
                files.append({
                    "name": filename,
                    "size": file_size,
                    "uploaded": file_time
                })
        
        return jsonify({"files": sorted(files, key=lambda x: x['uploaded'], reverse=True)})
        
    except Exception as e:
        logger.error(f"List Files Error: {str(e)}")
        return jsonify({"error": str(e)}), 500

@app.route('/ai_lab_chat', methods=['POST'])
def ai_lab_chat():
    """
    Analyze a file from AI Research Lab with AI
    Reads the stored file and answers questions about it
    """
    try:
        if not GENAI_AVAILABLE:
            return jsonify({
                "error": "AI service unavailable",
                "answer": "AI model is unavailable on current runtime. Set Python 3.11.9 in Render and redeploy."
            }), 503

        data = request.json
        filename = data.get('filename')
        query = data.get('query', '')
        
        logger.info(f"AI Lab Chat Request - Filename: {filename}, Query: {query[:100]}")
        
        if not filename:
            logger.error("No filename provided")
            return jsonify({"error": "No file specified"}), 400
        
        if not query:
            query = "Please analyze and summarize this document."
        
        # Get file from AI Lab storage
        LAB_FOLDER = os.path.join(UPLOAD_FOLDER, 'ai_lab_docs')
        file_path = os.path.join(LAB_FOLDER, secure_filename(filename))
        
        logger.info(f"Looking for file at: {file_path}")
        
        if not os.path.exists(file_path):
            logger.error(f"File not found at: {file_path}")
            return jsonify({"error": f"File not found: {filename}. Please ensure the file is fully uploaded."}), 404
        
        # Read file content based on file type
        file_content = ""
        try:
            if filename.lower().endswith('.pdf'):
                # Extract text from PDFī
                try:
                    pdf_document = fitz.open(file_path)
                    for page_num in range(len(pdf_document)):
                        page = pdf_document[page_num]
                        file_content += page.get_text() + "\n"
                    pdf_document.close()
                    logger.info(f"PDF extracted: {len(file_content)} characters")
                except Exception as pdf_error:
                    logger.error(f"PDF extraction error: {str(pdf_error)}")
                    file_content = f"[PDF file: {filename} - Could not extract text]"
                    
            elif filename.lower().endswith(('.txt', '.md', '.csv')):
                # Text files
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    file_content = f.read()
                logger.info(f"Text file read: {len(file_content)} characters")
                
            elif filename.lower().endswith(('.docx',)):
                # Word documents
                try:
                    from docx import Document
                    doc = Document(file_path)
                    file_content = '\n'.join([para.text for para in doc.paragraphs])
                    logger.info(f"DOCX extracted: {len(file_content)} characters")
                except:
                    file_content = "[DOCX file - extraction not available]"
            else:
                # Other file types
                try:
                    with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                        file_content = f.read()
                except:
                    file_content = f"[File type not directly supported: {filename}]"
                    
        except Exception as read_error:
            logger.error(f"File read error: {str(read_error)}")
            return jsonify({"error": f"Could not read file: {str(read_error)}"}), 500
        
        if not file_content or file_content.startswith('['):
            logger.warning(f"No readable content in file: {filename}")
            return jsonify({
                "error": "Could not extract content from file",
                "answer": f"The file '{filename}' appears to be empty or in an unsupported format. Please try with a text-based document (PDF, TXT, DOCX)."
            }), 400
        
        # Limit content to avoid token limits
        file_content = file_content[:8000]
        
        # Create prompt for AI analysis
        prompt = f"""You are an academic research assistant. A user has uploaded a document to analyze.

Document: {filename}

Document Content (first 8000 characters):
---
{file_content}
---

User Question/Request:
{query}

Please provide a helpful and detailed analysis or answer based on the document and the user's question."""
        
        logger.info(f"Sending to Gemini API - Content length: {len(file_content)}")
        
        # Generate response using Gemini
        try:
            model = genai.GenerativeModel(AI_MODEL_ID)
            response = model.generate_content(prompt)
            
            if response and response.text:
                logger.info(f"AI Response received - Length: {len(response.text)}")
                return jsonify({
                    "answer": response.text,
                    "filename": filename,
                    "status": "success"
                })
            else:
                logger.warning("AI returned empty response")
                return jsonify({
                    "error": "AI could not generate response",
                    "answer": "The AI model did not return a response. Please try again."
                }), 500
                
        except Exception as ai_error:
            logger.error(f"AI API Error: {str(ai_error)}")
            return jsonify({
                "error": f"AI Analysis Error: {str(ai_error)}",
                "answer": "The AI service encountered an error. Please check your API key and try again."
            }), 500
        
    except Exception as e:
        logger.error(f"AI Lab Chat Error: {str(e)}", exc_info=True)
        return jsonify({
            "error": f"Unexpected error: {str(e)}",
            "answer": "An unexpected error occurred. Check the logs for details."
        }), 500

@app.route('/delete-ai-lab-file', methods=['POST'])
def delete_ai_lab_file():
    """
    Delete a file from AI Research Lab
    """
    try:
        data = request.json
        filename = data.get('filename')
        
        if not filename:
            return jsonify({"error": "No file specified"}), 400
        
        LAB_FOLDER = os.path.join(UPLOAD_FOLDER, 'ai_lab_docs')
        file_path = os.path.join(LAB_FOLDER, secure_filename(filename))
        
        if not os.path.exists(file_path):
            return jsonify({"error": f"File not found: {filename}"}), 404
        
        os.remove(file_path)
        logger.info(f"File deleted from AI Lab: {filename}")
        
        return jsonify({
            "status": "success",
            "message": f"File '{filename}' deleted from AI Research Lab."
        })
        
    except Exception as e:
        logger.error(f"Delete File Error: {str(e)}")
        return jsonify({"error": str(e)}), 500
# --- ENHANCED AUTO-CLEANUP ---
def cleanup_loop():
    print("Cleanup thread started...")
    while True:
        now = time.time()
        count = 0
        try:
            for f in os.listdir(UPLOAD_FOLDER):
                p = os.path.join(UPLOAD_FOLDER, f)
                if f == 'student_data.json': continue
                # Skip folders (for example uploads/ai_lab_docs) to avoid permission errors.
                if not os.path.isfile(p):
                    continue
                if os.stat(p).st_mtime < now - 600:
                    try: 
                        os.remove(p)
                        count += 1
                    except Exception as e:
                        print(f"Error deleting {f}: {e}")
            if count > 0:
                print(f"Auto-Cleanup: Deleted {count} temporary files.")
        except Exception as e:
            print(f"Cleanup loop error: {e}")
        time.sleep(300)

if __name__ == '__main__':
    # Start cleanup thread
    if not app.debug or os.environ.get('WERKZEUG_RUN_MAIN') == 'true':
        Thread(target=cleanup_loop, daemon=True).start()
    
    # Run with appropriate settings based on environment
    if config.DEBUG:
        app.run(host='localhost', port=5000, debug=True)
    else:
        # For production, use gunicorn instead: gunicorn app:app
        logger.info("Running in production mode. Use 'gunicorn app:app' to start.")
        app.run(host='0.0.0.0', port=5000, debug=False)